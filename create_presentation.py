#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# Set slide size to widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define color palette
dark_charcoal = RGBColor(45, 45, 45)
accent_teal = RGBColor(0, 176, 176)
accent_gold = RGBColor(255, 193, 7)
accent_coral = RGBColor(255, 107, 107)
white = RGBColor(255, 255, 255)
light_gray = RGBColor(200, 200, 200)

def add_background(slide):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = dark_charcoal
    background.line.fill.background()
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

def style_text(text_frame, font_size=24, color=white, bold=False):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Segoe UI'
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold

def add_accent_bar(slide, color=accent_teal, height=0.1, y_position=0.5):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_position), Inches(12.3), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)
add_background(slide1)
add_accent_bar(slide1, accent_teal, 0.05, 1.5)
add_accent_bar(slide1, accent_coral, 0.05, 5.5)

title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
title_frame = title_box.text_frame
title_frame.text = 'THE CELLULAR-HORMONAL AXIS'
style_text(title_frame, 54, accent_teal, True)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = 'How Industrial Oils, Environmental Toxins, Chronic Stress\nand Circadian Disruption Create Modern Metabolic Disease'
style_text(subtitle_frame, 28, white)
for para in subtitle_frame.paragraphs:
    para.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = 'metabolic_dysfunction_presentation.pptx'
prs.save(output_path)
print(f'Basic presentation created: {output_path}')
