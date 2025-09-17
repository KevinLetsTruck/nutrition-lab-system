#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
dark_charcoal = RGBColor(45, 45, 45)
accent_teal = RGBColor(0, 176, 176)
accent_gold = RGBColor(255, 193, 7)
accent_coral = RGBColor(255, 107, 107)
white = RGBColor(255, 255, 255)
light_gray = RGBColor(200, 200, 200)

def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = dark_charcoal
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def style_text(text_frame, font_size=24, color=white, bold=False):
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.name = 'Segoe UI'
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold

def add_accent_bar(slide, color, height=0.05, y_pos=0.8):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.3), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

# Create slides
slides_data = [
    {
        'title': 'THE CELLULAR-HORMONAL AXIS',
        'subtitle': 'How Industrial Oils, Environmental Toxins, Chronic Stress\nand Circadian Disruption Create Modern Metabolic Disease',
        'footer': 'A Systems Approach to Understanding Metabolic Dysfunction',
        'color': accent_teal,
        'is_title': True
    },
    {
        'title': 'THE METABOLIC CRISIS',
        'content': [
            '• Over 1 billion people globally affected by metabolic dysfunction',
            '• Economic burden exceeding $583 billion annually in the US alone',
            '• Traditional "calories in, calories out" approach has failed',
            '• Dramatic increases coincide with industrial food transformation:',
            '   — Seed oil consumption: 0g (1865) → 80g daily (1999)',
            '   — Linoleic acid in body fat increased 136% since 1950s',
            '• Two interconnected root causes create a self-perpetuating cycle:',
            '   — Compromised cellular health',
            '   — Hormonal hijacking through stress and circadian disruption'
        ],
        'color': accent_coral
    },
    {
        'title': 'ROOT CAUSE 1: CELLULAR HEALTH COMPROMISE',
        'left_title': 'INDUSTRIAL SEED OILS',
        'left_content': [
            '• 60-80% of critical mitochondrial phospholipid (cardiolipin) displaced',
            '• 9-fold increase in inflammatory IL-8 secretion',
            '• Generates toxic 4-hydroxynonenal when heated',
            '• Japan case study (1960-2004):',
            '   — Seed oil: 9g → 39g daily',
            '   — Macular degeneration: 0.2% → 16.37% (82-fold increase)'
        ],
        'right_title': 'ENVIRONMENTAL TOXINS',
        'right_content': [
            '• 12 distinct mechanisms of metabolic disruption',
            '• Heavy metals: 23% increased metabolic syndrome risk per doubling of cadmium levels',
            '• BPA: 45% increased diabetes risk',
            '• Microplastics activate NLRP3 inflammasome pathways',
            '• Air pollution explains 17.9% of diabetes variation across US counties'
        ],
        'color': accent_teal,
        'two_column': True
    },
    {
        'title': 'ROOT CAUSE 2: HORMONAL HIJACKING',
        'left_title': 'CHRONIC STRESS ACTIVATION',
        'left_content': [
            '• HPA axis loses homeostatic control under sustained stress',
            '• Cortisol fails to return to baseline, creating new metabolic set point',
            '• Direct β-cell suppression + hepatic gluconeogenesis',
            '• Induces leptin resistance and ghrelin elevation',
            '• 30% increased diabetes incidence in shift workers'
        ],
        'right_title': 'CIRCADIAN DISRUPTION',
        'right_content': [
            '• 2-5 lux eliminates hormonal rhythms (nightlight level)',
            '• Urban light pollution: 150 lux vs. moonlight: 0.3 lux',
            '• Blue light (460-480nm) most potent for disruption',
            '• 2-hour tablet exposure reduces melatonin 55%, delays onset 1.5hrs',
            '• Clock disruption causes permanent insulin resistance in mice'
        ],
        'color': accent_gold,
        'two_column': True
    },
    {
        'title': 'SYNERGISTIC INTERACTIONS AMPLIFY DYSFUNCTION',
        'subtitle': 'THE VICIOUS CYCLE',
        'content': [
            '• Omega-6 oils generate pro-inflammatory compounds that amplify stress responses',
            '• Environmental toxins accumulate in fat tissue, releasing during altered circadian metabolism patterns',
            '• Oxidative stress from cellular damage disrupts the HPA axis',
            '• Stress hormones promote inflammatory gene expression, worsening cellular dysfunction',
            '• Sleep deprivation increases cortisol while decreasing insulin sensitivity',
            '• Korean data shows EXPONENTIAL (not linear) metabolic syndrome risk with multiple circadian-disrupting factors',
            '• Each disrupted system undermines the others, creating a metabolic prison increasingly difficult to escape'
        ],
        'color': accent_coral
    },
    {
        'title': 'CLINICAL EVIDENCE & VALIDATION',
        'subtitle': 'INTERVENTION STUDIES SHOW SUPERIOR OUTCOMES',
        'content': [
            '• Meta-analysis of 8 RCTs (2,839 participants): 2-fold greater metabolic syndrome resolution with comprehensive interventions',
            '• 28-day metabolic detoxification: Enhanced Phase II enzymes, improved antioxidant balance, reduced gamma-glutamyltransferase',
            '• Morning light therapy (>1000 lux): Improved insulin sensitivity',
            '• Amber-tinted glasses (80-100% blue light blocking): Better sleep quality and metabolic markers',
            '• Time-restricted eating aligned with circadian rhythms: Metabolic improvements in obese prediabetics',
            '',
            'GLOBAL IMPACT DATA',
            '• WHO: 26.1% of diabetes disability-adjusted life years attributable to air pollution',
            '• 1.55 million excess diabetes deaths globally from PM2.5 alone',
            '• $170+ billion annual obesity-related healthcare costs in US'
        ],
        'color': accent_teal
    },
    {
        'title': 'COMPREHENSIVE SOLUTION FRAMEWORK',
        'left_title': 'CELLULAR RESTORATION',
        'left_content': [
            '• Eliminate industrial seed oils (soybean, corn, canola, sunflower)',
            '• Prioritize omega-3 fatty acids and saturated fats',
            '• Support detoxification pathways:',
            '   — Phase I & II liver enzymes',
            '   — Glutathione production',
            '   — Methylation support',
            '• Reduce environmental toxin exposure:',
            '   — Filter water and air',
            '   — Choose organic foods',
            '   — Minimize plastic use'
        ],
        'right_title': 'HORMONAL OPTIMIZATION',
        'right_content': [
            '• Stress management protocols:',
            '   — HRV training',
            '   — Meditation/mindfulness',
            '   — Adaptogenic herbs',
            '• Circadian rhythm restoration:',
            '   — Morning bright light (>1000 lux)',
            '   — Blue light blocking after sunset',
            '   — Consistent sleep schedule',
            '   — Time-restricted eating',
            '• Environmental modifications:',
            '   — Dim lighting after dark',
            '   — Temperature regulation',
            '   — EMF reduction'
        ],
        'color': accent_gold,
        'two_column': True
    },
    {
        'title': 'KEY TAKEAWAYS',
        'content': [
            '1. Metabolic dysfunction stems from TWO interconnected root causes:',
            '   — Cellular health compromise from industrial oils and environmental toxins',
            '   — Hormonal hijacking through chronic stress and circadian disruption',
            '',
            '2. These factors create SYNERGISTIC effects greater than their sum, forming self-perpetuating cycles of metabolic damage',
            '',
            '3. Traditional "calories in, calories out" approaches fail because they ignore these foundational drivers',
            '',
            '4. Clinical evidence demonstrates that comprehensive interventions addressing both cellular and hormonal factors produce SUPERIOR outcomes',
            '',
            '5. Implementation requires systematic assessment and phased intervention targeting root causes, not just symptoms',
            '',
            '6. The framework explains individual variation in metabolic health and provides multiple intervention points beyond diet and exercise'
        ],
        'color': accent_teal
    },
    {
        'title': 'TRANSFORMING METABOLIC HEALTH',
        'subtitle': 'From Symptom Management\nto Root Cause Resolution',
        'footer': 'The future of metabolic medicine is systems-based, personalized, and preventive',
        'color': accent_teal,
        'is_final': True
    }
]

# Create slides
slide_layout = prs.slide_layouts[6]  # Blank layout

for i, slide_data in enumerate(slides_data):
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)
    add_accent_bar(slide, slide_data['color'])
    
    if slide_data.get('is_title'):
        # Title slide
        add_accent_bar(slide, accent_teal, 0.05, 1.5)
        add_accent_bar(slide, accent_coral, 0.05, 5.5)
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        style_text(title_frame, 54, accent_teal, True)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_data['subtitle']
        style_text(subtitle_frame, 28, white)
        for para in subtitle_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        
        author_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.8))
        author_frame = author_box.text_frame
        author_frame.text = slide_data['footer']
        style_text(author_frame, 18, light_gray)
        author_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    elif slide_data.get('is_final'):
        # Final slide
        add_accent_bar(slide, accent_teal, 0.03, 1.2)
        add_accent_bar(slide, accent_gold, 0.03, 1.3)
        add_accent_bar(slide, accent_coral, 0.03, 1.4)
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        style_text(title_frame, 48, accent_teal, True)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_data['subtitle']
        style_text(subtitle_frame, 32, white)
        for para in subtitle_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
        
        bottom_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(1))
        bottom_frame = bottom_box.text_frame
        bottom_frame.text = slide_data['footer']
        style_text(bottom_frame, 20, light_gray)
        bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    elif slide_data.get('two_column'):
        # Two column slide
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        style_text(title_frame, 32, slide_data['color'], True)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Left column
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.3))
        left_title_frame = left_title_box.text_frame
        left_title_frame.text = slide_data['left_title']
        style_text(left_title_frame, 20, slide_data['color'], True)
        
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.8))
        left_frame = left_box.text_frame
        left_frame.text = '\n'.join(slide_data['left_content'])
        style_text(left_frame, 16, white)
        
        # Right column
        right_title_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(0.3))
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = slide_data['right_title']
        style_text(right_title_frame, 20, slide_data['color'], True)
        
        right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.7), Inches(6), Inches(4.8))
        right_frame = right_box.text_frame
        right_frame.text = '\n'.join(slide_data['right_content'])
        style_text(right_frame, 16, white)
        
    else:
        # Regular content slide
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        style_text(title_frame, 32, slide_data['color'], True)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if slide_data.get('subtitle'):
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.1), Inches(11), Inches(0.4))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = slide_data['subtitle']
            style_text(subtitle_frame, 24, slide_data['color'], True)
            content_y = 1.6
        else:
            content_y = 1.3
        
        content_box = slide.shapes.add_textbox(Inches(1), Inches(content_y), Inches(11), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.text = '\n'.join(slide_data['content'])
        font_size = 18 if len('\n'.join(slide_data['content'])) > 800 else 20
        style_text(content_frame, font_size, white)
    
    print(f'Slide {i+1} created')

# Save presentation
prs.save('metabolic_dysfunction_presentation.pptx')
print('Full presentation created successfully!')
